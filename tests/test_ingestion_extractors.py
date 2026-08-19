import docx as docx_lib
from pptx import Presentation

import course_mcp_server.ingestion as ingestion

from course_mcp_server.ingestion import extract_source


def _build_docx(path, *, with_table: bool = False):
    document = docx_lib.Document()
    document.add_paragraph("Safety Policy", style="Heading 1")
    document.add_paragraph("Inspect equipment before use.")
    if with_table:
        table = document.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "Item"
        table.rows[0].cells[1].text = "Status"
        table.rows[1].cells[0].text = "Helmet"
        table.rows[1].cells[1].text = "OK"
    document.save(str(path))


def test_docx_extraction_reads_headings_and_body_text(tmp_path):
    docx = tmp_path / "policy.docx"
    _build_docx(docx)

    result = extract_source(docx, "docx")

    assert "Safety Policy" in result.text
    assert "Inspect equipment" in result.text
    assert result.headings == ["Safety Policy"]
    assert result.references == ["section:Safety Policy"]
    assert result.tables == []


def test_docx_extraction_reads_table_cells(tmp_path):
    docx = tmp_path / "policy-with-table.docx"
    _build_docx(docx, with_table=True)

    result = extract_source(docx, "docx")

    assert result.tables == [["Item | Status", "Helmet | OK"]]
    assert "Item | Status" in result.text
    assert "Helmet | OK" in result.text


def test_pptx_extraction_reads_slide_order_and_notes(tmp_path):
    pptx = tmp_path / "deck.pptx"
    presentation = Presentation()
    layout = presentation.slide_layouts[1]

    slide1 = presentation.slides.add_slide(layout)
    slide1.shapes.title.text = "Slide Title"
    slide1.notes_slide.notes_text_frame.text = "Speaker note"

    slide2 = presentation.slides.add_slide(layout)
    slide2.shapes.title.text = "Second Slide"

    presentation.save(str(pptx))

    result = extract_source(pptx, "pptx")

    assert "Slide 1: Slide Title" in result.text
    assert "Slide 2: Second Slide" in result.text
    assert "Notes 1: Speaker note" in result.text
    assert result.references == ["slide:1", "slide:2", "notes:1"]


def test_youtube_transcript_import_requires_controlled_text_file(tmp_path):
    transcript = tmp_path / "video.txt"
    transcript.write_text("00:00 Welcome\n00:12 First safety rule", encoding="utf-8")

    result = extract_source(transcript, "youtube")

    assert "First safety rule" in result.text
    assert result.references == ["timestamp:00:00", "timestamp:00:12"]


def test_pdf_extraction_reports_page_markers_from_text_pdf(tmp_path):
    pdf = tmp_path / "sample.pdf"
    pdf.write_bytes(b"%PDF-1.4\nBT (Page One) Tj ET\n%%Page: 2 2\nBT (Page Two) Tj ET")

    result = extract_source(pdf, "pdf")

    assert "Page One" in result.text
    assert "Page Two" in result.text
    assert "page:1" in result.references


def test_pdf_extraction_uses_page_aware_pypdf(monkeypatch, tmp_path):
    pdf = tmp_path / "thirty-pages.pdf"
    pdf.write_bytes(b"%PDF-1.4 test fixture")

    class Page:
        def __init__(self, number):
            self.number = number

        def extract_text(self):
            return f"Policy page {self.number} — café safety"

    class Reader:
        def __init__(self, _path):
            self.pages = [Page(number) for number in range(1, 31)]

    import pypdf

    monkeypatch.setattr(pypdf, "PdfReader", Reader)
    result = ingestion.extract_source(pdf, "pdf")

    assert "[page 1]" in result.text
    assert "[page 30]" in result.text
    assert "café safety" in result.text
    assert result.references == [f"page:{number}" for number in range(1, 31)]
    assert result.warnings == []
