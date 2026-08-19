from __future__ import annotations

import re
from dataclasses import dataclass, field
from html import unescape
from pathlib import Path


@dataclass(frozen=True)
class ExtractedSource:
    text: str
    headings: list[str] = field(default_factory=list)
    references: list[str] = field(default_factory=list)
    tables: list[list[str]] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


def _collapse(text: str) -> str:
    return re.sub(r"\s+", " ", unescape(text)).strip()


def _extract_docx(path: Path) -> ExtractedSource:
    import docx  # type: ignore

    document = docx.Document(str(path))
    paragraphs: list[str] = []
    headings: list[str] = []
    for paragraph in document.paragraphs:
        text = _collapse(paragraph.text or "")
        if not text:
            continue
        paragraphs.append(text)
        style = paragraph.style
        style_name = getattr(style, "name", "") or ""
        if "Heading" in style_name:
            headings.append(text)

    tables: list[list[str]] = []
    table_lines: list[str] = []
    for table_index, table in enumerate(document.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            cells = [_collapse(cell.text or "") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            tables.append(rows)
            table_lines.append(f"[table {table_index}]")
            table_lines.extend(rows)

    refs = [f"section:{heading}" for heading in headings] or ["section:body"]
    return ExtractedSource(
        text="\n".join(paragraphs + table_lines),
        headings=headings,
        references=refs,
        tables=tables,
    )


def _pptx_shape_text(shape: object) -> str:
    if getattr(shape, "has_text_frame", False):
        return _collapse(shape.text_frame.text or "")  # type: ignore[attr-defined]
    if getattr(shape, "has_table", False):
        cells: list[str] = []
        for row in shape.table.rows:  # type: ignore[attr-defined]
            for cell in row.cells:
                text = _collapse(cell.text or "")
                if text:
                    cells.append(text)
        return " ".join(cells)
    return ""


def _extract_pptx(path: Path) -> ExtractedSource:
    from pptx import Presentation  # type: ignore

    presentation = Presentation(str(path))
    slides = list(presentation.slides)
    lines: list[str] = []
    refs: list[str] = []

    for slide_number, slide in enumerate(slides, start=1):
        parts = [_pptx_shape_text(shape) for shape in slide.shapes]
        text = " ".join(part for part in parts if part)
        if text:
            lines.append(f"Slide {slide_number}: {text}")
            refs.append(f"slide:{slide_number}")

    for slide_number, slide in enumerate(slides, start=1):
        if slide.has_notes_slide:
            notes = _collapse(slide.notes_slide.notes_text_frame.text or "")
            if notes:
                lines.append(f"Notes {slide_number}: {notes}")
                refs.append(f"notes:{slide_number}")

    return ExtractedSource(text="\n".join(lines), references=refs)


def _extract_pdf(path: Path) -> ExtractedSource:
    warnings: list[str] = []
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        pages: list[str] = []
        references: list[str] = []
        for page_number, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            references.append(f"page:{page_number}")
            if page_text:
                pages.append(f"[page {page_number}]\n{page_text}")
            else:
                warnings.append(f"Page {page_number} has no extractable text")
        return ExtractedSource(
            text="\n\n".join(pages),
            references=references or ["page:1"],
            warnings=warnings,
        )
    except Exception as exc:  # malformed PDFs and optional dependency fallback
        warnings.append(f"pypdf extraction failed; used compatibility fallback ({type(exc).__name__})")

    raw = path.read_bytes().decode("latin-1", errors="ignore")
    strings = re.findall(r"\(([^()] {0,0}.*?)\)\s*Tj", raw)
    if not strings:
        strings = re.findall(r"[A-Za-z][A-Za-z0-9 ,.;:!?'-]{4,}", raw)
    text = "\n".join(_collapse(item) for item in strings if _collapse(item))
    page_count = max(1, raw.count("/Page") or raw.count("%%Page") or 1)
    return ExtractedSource(
        text=text,
        references=[f"page:{index}" for index in range(1, page_count + 1)],
        warnings=warnings,
    )


def _extract_timed_text(path: Path) -> ExtractedSource:
    text = path.read_text(encoding="utf-8", errors="ignore")
    refs = re.findall(r"\b(\d{1,2}:\d{2}(?::\d{2})?)\b", text)
    return ExtractedSource(text=text, references=[f"timestamp:{stamp}" for stamp in refs])


def _extract_raw_text(path: Path) -> ExtractedSource:
    text = path.read_text(encoding="utf-8", errors="ignore")
    headings = [
        line.lstrip("#").strip()
        for line in text.splitlines()
        if line.strip().startswith("#") and line.lstrip("#").strip()
    ]
    refs = [f"section:{heading}" for heading in headings] or ["line:1"]
    return ExtractedSource(text=text, headings=headings, references=refs)


def extract_source(path: Path | str, source_type: str) -> ExtractedSource:
    source_path = Path(path)
    if source_type == "docx":
        return _extract_docx(source_path)
    if source_type in {"pptx", "ppt"}:
        return _extract_pptx(source_path)
    if source_type == "pdf":
        return _extract_pdf(source_path)
    if source_type in {"youtube", "website"}:
        return _extract_timed_text(source_path)
    return _extract_raw_text(source_path)
