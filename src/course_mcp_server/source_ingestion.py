from __future__ import annotations

import hashlib
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Any

SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx"}
MAX_SOURCE_BYTES = int(os.getenv("MAX_SOURCE_BYTES", str(25 * 1024 * 1024)))


class SourceIngestionError(RuntimeError):
    pass


@dataclass(frozen=True)
class ExtractedSource:
    source_id: str
    source_type: str
    title: str
    text: str
    chunks: list[dict[str, Any]]
    sha256: str
    warnings: list[str]


def resolve_under_root(root: Path | str, relative_name: str) -> Path:
    root_path = Path(root).resolve()
    candidate = (root_path / relative_name).resolve()
    if root_path != candidate and root_path not in candidate.parents:
        raise SourceIngestionError("Source path escapes the allowed upload root")
    return candidate


def _hash_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _chunk_text(text: str, *, source_id: str, source_type: str, chunk_size: int = 2500) -> list[dict[str, Any]]:
    normalized = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    chunks: list[dict[str, Any]] = []
    start = 0
    index = 1
    while start < len(normalized):
        part = normalized[start : start + chunk_size]
        chunks.append({
            "chunk_id": f"{source_id}_chunk_{index}",
            "source_id": source_id,
            "source_type": source_type,
            "text": part,
            "char_start": start,
            "char_end": start + len(part),
        })
        start += chunk_size
        index += 1
    return chunks


def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _extract_pdf(path: Path) -> tuple[str, list[str]]:
    warnings: list[str] = []
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency
        raise SourceIngestionError("pypdf is required for PDF ingestion") from exc
    reader = PdfReader(str(path))
    pages: list[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if not text.strip():
            warnings.append(f"Page {page_number} has no extractable text")
        pages.append(f"[page {page_number}]\n{text}")
    return "\n\n".join(pages), warnings


def _extract_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency
        raise SourceIngestionError("python-docx is required for DOCX ingestion") from exc
    document = docx.Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table_index, table in enumerate(document.tables, start=1):
        parts.append(f"[table {table_index}]")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency
        raise SourceIngestionError("python-pptx is required for PPTX ingestion") from exc
    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[slide {slide_number}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[slide {slide_number} notes]\n{notes}")
    return "\n".join(parts)


def ingest_source_from_upload(upload_root: Path | str, upload_id: str) -> ExtractedSource:
    """Extract text from an uploaded file stored under upload_root.

    The caller must pass a safe upload ID such as `course-sop.pdf`, not an absolute path.
    This function rejects path traversal and size abuse before parsing.
    """
    if not upload_id or upload_id.startswith(("/", "\\")) or ".." in Path(upload_id).parts:
        raise SourceIngestionError("Invalid upload_id")

    path = resolve_under_root(upload_root, upload_id)
    if not path.exists() or not path.is_file():
        raise SourceIngestionError("Uploaded source was not found")

    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise SourceIngestionError(f"Unsupported source type: {suffix}")

    size = path.stat().st_size
    if size > MAX_SOURCE_BYTES:
        raise SourceIngestionError("Uploaded source exceeds maximum allowed size")

    raw = path.read_bytes()
    sha = _hash_bytes(raw)
    source_id = f"src_{sha[:16]}"
    warnings: list[str] = []

    if suffix in {".txt", ".md"}:
        text = _extract_txt(path)
    elif suffix == ".pdf":
        text, warnings = _extract_pdf(path)
    elif suffix == ".docx":
        text = _extract_docx(path)
    elif suffix == ".pptx":
        text = _extract_pptx(path)
    else:  # pragma: no cover
        raise SourceIngestionError(f"Unsupported source type: {suffix}")

    if not text.strip():
        warnings.append("No useful text could be extracted from source")

    source_type = suffix.lstrip(".")
    return ExtractedSource(
        source_id=source_id,
        source_type=source_type,
        title=path.name,
        text=text,
        chunks=_chunk_text(text, source_id=source_id, source_type=source_type),
        sha256=sha,
        warnings=warnings,
    )
